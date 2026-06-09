"""
모바일 프로그래밍 최종 발표 PPT 생성.
- 강의: 모바일 프로그래밍
- 팀: 봉준표(팀장), 이동제, 윤준영
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

OUT = "/Users/bongjunpyo/smartfarm-upgrade"
A = f"{OUT}/report_assets"
FILE = f"{OUT}/모바일프로그래밍_최종발표.pptx"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

# 색상
GREEN  = RGBColor(0x16, 0xa3, 0x4a)
DGREEN = RGBColor(0x15, 0x80, 0x3d)
LGREEN = RGBColor(0xdc, 0xfc, 0xe7)
RED    = RGBColor(0xdc, 0x26, 0x26)
LRED   = RGBColor(0xfe, 0xe2, 0xe2)
NAVY   = RGBColor(0x0f, 0x17, 0x2a)
SLATE  = RGBColor(0x47, 0x55, 0x69)
MUTED  = RGBColor(0x64, 0x74, 0x8b)
LIGHT  = RGBColor(0xf1, 0xf5, 0xf9)
BORDER = RGBColor(0xcb, 0xd5, 0xe1)
WHITE  = RGBColor(0xff, 0xff, 0xff)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
ORANGE = RGBColor(0xea, 0x58, 0x0c)
YELLOW = RGBColor(0xf5, 0x9e, 0x0b)

FONT  = "AppleGothic"

# ============== helpers ==============
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def textbox(slide, x, y, w, h, text, size=18, bold=False, color=NAVY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def rounded(slide, x, y, w, h, fill=WHITE, line=BORDER, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.08
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh

def page_header(slide, title, subtitle=None, accent=GREEN):
    rect(slide, 0, 0, SW, Inches(0.9), accent)
    textbox(slide, Inches(0.5), Inches(0.18), Inches(11), Inches(0.55),
            title, size=26, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, Inches(0.55), Inches(0.55), Inches(11), Inches(0.4),
                subtitle, size=12, color=WHITE)
    # 푸터
    textbox(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
            "모바일 프로그래밍 — 비닐하우스 질병 감지 AI 앱", size=10, color=MUTED)
    textbox(slide, Inches(11.3), Inches(7.05), Inches(2), Inches(0.3),
            "봉준표 · 이동제 · 윤준영", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_image_fit(slide, path, x, y, w, h):
    """이미지 비율 유지하면서 영역 안에 맞춰 넣기"""
    from PIL import Image
    img = Image.open(path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        nw = w; nh = int(w / img_ratio)
    else:
        nh = h; nw = int(h * img_ratio)
    nx = x + (w - nw) // 2
    ny = y + (h - nh) // 2
    slide.shapes.add_picture(path, nx, ny, nw, nh)

# ============== 슬라이드 1: 표지 ==============
s = add_blank_slide()
set_bg(s, NAVY)
# 데코 사각형
deco1 = rect(s, 0, Inches(5.2), SW, Inches(2.3), DGREEN)
deco2 = rect(s, 0, Inches(6.7), SW, Inches(0.8), GREEN)
textbox(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
        "2026학년도 1학기  /  모바일 프로그래밍", size=14, color=LIGHT)
textbox(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.4),
        "최종 발표 자료", size=12, color=RGBColor(0x94, 0xa3, 0xb8))
textbox(s, Inches(0.7), Inches(1.7), Inches(12.2), Inches(1.6),
        "스마트팜 모니터링 안드로이드 앱", size=54, bold=True, color=WHITE)
textbox(s, Inches(0.7), Inches(3.0), Inches(12.2), Inches(0.7),
        "SmartFarm — AI 진단을 활용한 비닐하우스 구역 모니터링 모바일 앱",
        size=18, color=RGBColor(0xa7, 0xf3, 0xd0))
# 팀 카드
rounded(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.95),
        fill=RGBColor(0x1e, 0x29, 0x3b), line=None)
textbox(s, Inches(1.0), Inches(4.10), Inches(2.4), Inches(0.35),
        "팀장", size=12, color=GREEN, bold=True)
textbox(s, Inches(1.0), Inches(4.40), Inches(3.5), Inches(0.5),
        "봉  준  표", size=22, bold=True, color=WHITE)
textbox(s, Inches(5.0), Inches(4.10), Inches(2.4), Inches(0.35),
        "팀원", size=12, color=GREEN, bold=True)
textbox(s, Inches(5.0), Inches(4.40), Inches(3.5), Inches(0.5),
        "이  동  제", size=22, bold=True, color=WHITE)
textbox(s, Inches(8.8), Inches(4.10), Inches(2.4), Inches(0.35),
        "팀원", size=12, color=GREEN, bold=True)
textbox(s, Inches(8.8), Inches(4.40), Inches(3.5), Inches(0.5),
        "윤  준  영", size=22, bold=True, color=WHITE)
# 하단
textbox(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.45),
        "Android (Java)  ·  Material Design  ·  CameraX  ·  AI 진단 연동 (ONNX)",
        size=16, color=WHITE)
textbox(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
        "발표일: 2026-06-07", size=12, color=WHITE, align=PP_ALIGN.RIGHT)

# ============== 슬라이드 2: 목차 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "목차", "Table of Contents")
items = [
    ("01", "프로젝트 개요",            "강의 / 팀 / 앱 핵심 지표"),
    ("02", "해결하고자 하는 문제",     "농가 페인 포인트와 우리의 해결책"),
    ("03", "주요 기능 (앱 중심)",      "구역 현황 · 사진 기록 · 알림 · 시각화"),
    ("04", "시스템 아키텍처",          "Android 앱 ↔ FastAPI ↔ AI 모델"),
    ("05", "앱에 탑재된 AI (요약)",    "모델 학습 결과와 앱에서의 활용 방식"),
    ("06", "AI 응답 → 앱 시각화",       "분류 결과 + 그리드 오버레이 표시"),
    ("07", "안드로이드 앱 화면",        "실제 구동 캡처 4종"),
    ("08", "앱 사용 시나리오",         "위험 감지 → 상세 확인 → 알림 처리"),
    ("09", "역할 분담",                "팀원별 구현 부분"),
    ("10", "개발 과정에서 느낀 점",     "모바일 앱 개발 회고"),
]
col1_x = Inches(0.7); col2_x = Inches(6.85)
for i, (n, k, v) in enumerate(items):
    col = 0 if i < 5 else 1
    row = i if col == 0 else i - 5
    x = col1_x if col == 0 else col2_x
    y = Inches(1.25 + row * 1.08)
    rounded(s, x, y, Inches(5.9), Inches(0.95), fill=LIGHT, line=BORDER)
    rect(s, x, y, Inches(0.18), Inches(0.95), GREEN)
    textbox(s, x + Inches(0.4), y + Inches(0.08), Inches(0.9), Inches(0.4),
            n, size=22, bold=True, color=GREEN)
    textbox(s, x + Inches(1.3), y + Inches(0.05), Inches(4.5), Inches(0.45),
            k, size=16, bold=True, color=NAVY)
    textbox(s, x + Inches(1.3), y + Inches(0.45), Inches(4.5), Inches(0.45),
            v, size=11, color=MUTED)

# ============== 슬라이드 3: 프로젝트 개요 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "01. 프로젝트 개요", "스마트팜 모니터링 안드로이드 앱")
# 한 줄 요약
rounded(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.95), fill=LGREEN, line=GREEN)
textbox(s, Inches(0.95), Inches(1.30), Inches(11.6), Inches(0.5),
        "농가가 스마트폰 한 장으로 비닐하우스 전 구역의 병해를 모니터링하는 안드로이드 앱",
        size=18, bold=True, color=DGREEN)
textbox(s, Inches(0.95), Inches(1.65), Inches(11.6), Inches(0.4),
        "구역 현황 · 사진 기록 · 실시간 알림 · AI 진단 결과 시각화를 한 화면에서 제공하는 모바일 통합 솔루션",
        size=12, color=DGREEN)

# 정보 카드 4개
cards = [
    ("강의명",         "모바일 프로그래밍",                          GREEN),
    ("개발 기간",       "2026.05.01 ~ 2026.06.07 (약 5주)",          BLUE),
    ("개발 환경",       "Android Studio (Java) + Python (PyTorch)",  ORANGE),
    ("최종 목표",       "온디바이스 식물 병해 진단 + 구역별 모니터링", DGREEN),
]
cx = Inches(0.7); cy = Inches(2.4); cw = Inches(2.95); ch = Inches(1.4)
for i, (k, v, c) in enumerate(cards):
    x = cx + i * Inches(3.05)
    rounded(s, x, cy, cw, ch, fill=WHITE, line=BORDER)
    rect(s, x, cy, cw, Inches(0.08), c)
    textbox(s, x + Inches(0.2), cy + Inches(0.20), cw - Inches(0.3), Inches(0.4),
            k, size=12, bold=True, color=c)
    textbox(s, x + Inches(0.2), cy + Inches(0.55), cw - Inches(0.3), Inches(0.85),
            v, size=14, color=NAVY)

# 핵심 지표
y2 = Inches(4.05)
textbox(s, Inches(0.7), y2, Inches(8), Inches(0.4),
        "앱 핵심 지표", size=18, bold=True, color=NAVY)
metrics = [
    ("4 화면",  "주요 Activity",    GREEN),
    ("20 구역", "지원 카메라 수",    BLUE),
    ("1~2초",  "AI 진단 응답",      ORANGE),
    ("23종",   "감지 가능 병해",     DGREEN),
]
my = Inches(4.55); mw = Inches(2.95); mh = Inches(2.0)
for i, (v, k, c) in enumerate(metrics):
    x = Inches(0.7) + i * Inches(3.05)
    rounded(s, x, my, mw, mh, fill=LIGHT, line=BORDER)
    textbox(s, x, my + Inches(0.35), mw, Inches(0.9),
            v, size=42, bold=True, color=c, align=PP_ALIGN.CENTER)
    textbox(s, x, my + Inches(1.3), mw, Inches(0.5),
            k, size=14, color=SLATE, align=PP_ALIGN.CENTER)

# ============== 슬라이드 4: 해결하고자 하는 문제 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "02. 해결하고자 하는 사용자의 문제", "Why we built this")
# 좌측: 문제
rounded(s, Inches(0.7), Inches(1.2), Inches(6.0), Inches(5.6), fill=LRED, line=RED)
textbox(s, Inches(0.95), Inches(1.4), Inches(5.5), Inches(0.5),
        "현재 농가가 겪는 문제", size=18, bold=True, color=RED)
problems = [
    ("초기 병해 발견이 어렵다",
     "잎 한 장의 갈변·반점을 육안으로 식별하기 힘들어 발견이 늦어진다."),
    ("전문가 진단 비용이 비싸다",
     "농촌진흥청·식물병원 의뢰는 시간/비용 부담이 크고 현장 출장이 어렵다."),
    ("비닐하우스 환경 — 인터넷 불안정",
     "외곽 농장은 LTE/Wi-Fi 신호가 약해 클라우드 API에 의존할 수 없다."),
    ("구역별 모니터링 부재",
     "5~10평 구역 단위로 어디서 병이 시작됐는지 파악할 도구가 없다."),
    ("작물·질병 종류가 다양하다",
     "토마토·감자·사과·옥수수·피망 등 다품종 농가가 각각 별도 도구를 써야 한다."),
]
py = Inches(1.95)
for i, (t, d) in enumerate(problems):
    rect(s, Inches(0.95), py + Inches(0.05), Inches(0.05), Inches(0.85), RED)
    textbox(s, Inches(1.1), py, Inches(5.4), Inches(0.45),
            f"{i+1}. {t}", size=14, bold=True, color=NAVY)
    textbox(s, Inches(1.1), py + Inches(0.40), Inches(5.4), Inches(0.55),
            d, size=11, color=SLATE)
    py += Inches(0.95)

# 우측: 해결책
rounded(s, Inches(6.95), Inches(1.2), Inches(6.0), Inches(5.6), fill=LGREEN, line=GREEN)
textbox(s, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
        "우리의 해결책", size=18, bold=True, color=DGREEN)
solutions = [
    ("스마트폰 카메라 한 장으로 진단",
     "별도 장비 없이 잎 사진만 찍으면 23종 질병을 즉시 판별."),
    ("학습된 EfficientNet-B3 (Val 95.3%)",
     "ImageNet pretrained + 20,000장 학습으로 전문가급 정확도."),
    ("온디바이스 추론 (ONNX Runtime)",
     "인터넷 불필요 — 모델이 단말기에 내장되어 오프라인 동작."),
    ("10×10 그리드 오버레이",
     "병반 분포를 화면 위에 시각화해 농가가 직접 구역별 진단 가능."),
    ("23 클래스 다품종 지원",
     "감자/토마토/피망/사과/옥수수까지 한 앱으로 통합 진단."),
]
sy = Inches(1.95)
for i, (t, d) in enumerate(solutions):
    rect(s, Inches(7.2), sy + Inches(0.05), Inches(0.05), Inches(0.85), GREEN)
    textbox(s, Inches(7.35), sy, Inches(5.4), Inches(0.45),
            f"{chr(0x2713)}  {t}", size=14, bold=True, color=DGREEN)
    textbox(s, Inches(7.35), sy + Inches(0.40), Inches(5.4), Inches(0.55),
            d, size=11, color=SLATE)
    sy += Inches(0.95)

# ============== 슬라이드 5: 주요 기능 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "03. 주요 기능", "Key Features")
features = [
    ("🌿",  "구역 현황 대시보드",   "정상·위험·주의 카운트와 카드뷰로 전 구역 상태 한눈에",     GREEN),
    ("🖼",  "사진 기록 갤러리",     "날짜·구역·질병 태그 필터로 누적된 진단 사진 조회",        BLUE),
    ("🔍",  "구역 상세 모달",       "선택 구역의 진단 결과 + 신뢰도 + 촬영 시각까지 표시",    ORANGE),
    ("🔔",  "실시간 알림 탭",       "위험 감지 시 자동 누적 — 미확인 알림 뱃지 표시",         RED),
    ("📸",  "이미지 업로드",        "FastAPI /analyze로 multipart 전송, AI 응답 즉시 반영",   DGREEN),
    ("🎨",  "Material Design UI",  "그린 톤 브랜드 컬러 · 카드뷰 · 하단 탭 네비게이션",       YELLOW),
    ("🗺",  "병반 위치 오버레이",   "AI가 보내준 그리드 결과를 빨강/초록 시각화로 출력",      BLUE),
    ("⚙",  "구역별 필터 칩",       "전체/A1~B4 등 카메라별 필터로 데이터 분리 조회",         DGREEN),
]
cx0 = Inches(0.7); cy0 = Inches(1.15); cw = Inches(3.05); ch = Inches(1.42)
for i, (icon, k, v, c) in enumerate(features):
    col = i % 4; row = i // 4
    x = cx0 + col * Inches(3.15)
    y = cy0 + row * Inches(1.50)
    rounded(s, x, y, cw, ch, fill=WHITE, line=BORDER)
    rect(s, x, y, Inches(0.10), ch, c)
    textbox(s, x + Inches(0.25), y + Inches(0.1), Inches(0.6), Inches(0.6),
            icon, size=26, color=c)
    textbox(s, x + Inches(0.95), y + Inches(0.18), cw - Inches(1.0), Inches(0.4),
            k, size=14, bold=True, color=NAVY)
    textbox(s, x + Inches(0.25), y + Inches(0.75), cw - Inches(0.4), Inches(0.6),
            v, size=10, color=SLATE)

# 하단 강조
y3 = Inches(4.30)
rounded(s, Inches(0.7), y3, Inches(12.0), Inches(2.45), fill=NAVY, line=None)
textbox(s, Inches(1.0), y3 + Inches(0.18), Inches(11), Inches(0.45),
        "왜 \"모바일 앱\"이어야 했는가", size=18, bold=True, color=WHITE)
points = [
    ("현장성",   "농가가 비닐하우스 안에서 손에 들고 바로 사용 — PC/웹 대비 압도적 접근성"),
    ("실시간",   "구역 상태 변화·위험 알림이 푸시로 즉시 전달되어 대응 지연 최소화"),
    ("올인원",   "촬영 · 결과 확인 · 이력 조회 · 알림을 한 앱에서 끝내는 통합 UX"),
    ("확장성",   "Android 단일 코드로 농가 보급형부터 관제 태블릿까지 폼팩터 대응 가능"),
]
py = y3 + Inches(0.75);
for i, (k, v) in enumerate(points):
    col = i % 2; row = i // 2
    x = Inches(1.0) + col * Inches(5.85)
    y = py + row * Inches(0.78)
    rect(s, x, y + Inches(0.08), Inches(0.06), Inches(0.5), GREEN)
    textbox(s, x + Inches(0.15), y, Inches(2.0), Inches(0.4),
            k, size=14, bold=True, color=GREEN)
    textbox(s, x + Inches(2.0), y + Inches(0.04), Inches(3.8), Inches(0.6),
            v, size=11, color=WHITE)

# ============== 슬라이드 6: 시스템 아키텍처 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "04. 시스템 아키텍처", "Android × AI × Backend")
add_image_fit(s, f"{A}/architecture.png",
              Inches(0.5), Inches(1.05), Inches(12.3), Inches(5.4))
textbox(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
        "Android 앱 — 학습된 ONNX 모델을 assets에 내장 — 백엔드는 분석 기록/알림 담당",
        size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ============== 슬라이드 7: 앱에 탑재된 AI (요약) ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "05. 앱에 탑재된 AI 모델 (요약)",
            "복잡한 학습 과정은 서버에서, 앱은 결과만 받아 표시")

# 왼쪽 상단 — 학습 곡선 작게
add_image_fit(s, f"{A}/training_curve.png",
              Inches(0.4), Inches(1.05), Inches(7.5), Inches(3.4))
# 우측 — 핵심 결과 카드 4개
rs = [
    ("95.30%",  "Val Accuracy",     GREEN),
    ("23종",    "감지 가능 병해",     BLUE),
    ("~41 MB",  "ONNX 모델 크기",     ORANGE),
    ("1~2초",   "백엔드 추론 응답",   DGREEN),
]
for i, (v, k, c) in enumerate(rs):
    col = i % 2; row = i // 2
    x = Inches(8.2) + col * Inches(2.45)
    y = Inches(1.05) + row * Inches(1.75)
    rounded(s, x, y, Inches(2.35), Inches(1.6), fill=WHITE, line=c, line_w=2)
    textbox(s, x, y + Inches(0.20), Inches(2.35), Inches(0.6),
            v, size=24, bold=True, color=c, align=PP_ALIGN.CENTER)
    textbox(s, x, y + Inches(0.95), Inches(2.35), Inches(0.5),
            k, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# 하단 — 앱 관점 설명
rounded(s, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.0), fill=NAVY, line=None)
textbox(s, Inches(0.7), Inches(4.85), Inches(11.5), Inches(0.5),
        "앱은 AI를 어떻게 사용하는가?", size=16, bold=True, color=WHITE)
points = [
    ("📷",  "사용자가 잎 사진 촬영 → 앱은 multipart로 백엔드 /analyze 호출"),
    ("🧠",  "백엔드가 ONNX 모델로 추론 → disease_type + confidence 반환"),
    ("🗺",  "앱은 결과를 받아 카드뷰의 색상·라벨·신뢰도로 즉시 시각화"),
    ("🔔",  "위험 임계치 초과 시 알림 탭에 자동 누적 — 별도 동작 불필요"),
]
yy = Inches(5.35)
for i, (ico, txt) in enumerate(points):
    col = i % 2; row = i // 2
    x = Inches(0.7) + col * Inches(6.1)
    y = yy + row * Inches(0.65)
    textbox(s, x, y, Inches(0.5), Inches(0.5),
            ico, size=18, color=GREEN)
    textbox(s, x + Inches(0.55), y + Inches(0.05), Inches(5.5), Inches(0.6),
            txt, size=11, color=WHITE)

# ============== 슬라이드 9: 데모 - 추론 갤러리 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "06. 앱이 받는 AI 응답 — 분류 결과 예시",
            "백엔드가 반환하는 disease_type + confidence를 앱이 카드로 시각화")
add_image_fit(s, f"{A}/inference_gallery.png",
              Inches(0.3), Inches(1.05), Inches(12.7), Inches(5.8))
textbox(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
        "각 이미지 위에는 모델이 예측한 클래스명과 신뢰도(%)가 표기됨",
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ============== 슬라이드 10: 데모 - 그리드 오버레이 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "06. 앱 안의 \"병반 위치 시각화\"",
            "AI가 반환한 셀별 결과를 앱이 빨강/초록 그리드로 화면에 그려줌")
add_image_fit(s, f"{A}/grid_overlay.png",
              Inches(0.5), Inches(1.05), Inches(8.5), Inches(5.4))
# 우측 설명
rounded(s, Inches(9.3), Inches(1.2), Inches(3.7), Inches(5.2), fill=LIGHT, line=BORDER)
textbox(s, Inches(9.5), Inches(1.35), Inches(3.4), Inches(0.4),
        "그리드 추론 결과", size=15, bold=True, color=NAVY)
infos = [
    ("입력",     "TomatoEarlyBlight1.JPG"),
    ("그리드",   "10 × 10 = 100 셀"),
    ("병해 셀",  "49 / 100"),
    ("정상 셀",  "51 / 100"),
    ("병해 비율", "49 %"),
    ("최종 판정", "Early Blight\n— 즉시 방제 권장"),
]
yi = Inches(1.85)
for k, v in infos:
    textbox(s, Inches(9.5), yi, Inches(3.4), Inches(0.3),
            k, size=11, color=MUTED)
    textbox(s, Inches(9.5), yi + Inches(0.28), Inches(3.4), Inches(0.5),
            v, size=13, bold=True, color=NAVY)
    yi += Inches(0.65)

# ============== 슬라이드 11: 안드로이드 UI — 실제 캡처 ==============
PIC = "/Users/bongjunpyo/Desktop/app_picture"

s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "07. 안드로이드 앱 화면 — 실제 구동 캡처",
            "Splash → 구역 현황 → 사진 기록 → 알림")
shots = [
    (f"{PIC}/KakaoTalk_Photo_2026-06-08-12-09-08.png", "① 스플래시",     "앱 진입 / 브랜딩"),
    (f"{PIC}/KakaoTalk_Photo_2026-06-08-12-08-33.png", "② 구역 현황",     "정상18 위험2 / 카드뷰"),
    (f"{PIC}/KakaoTalk_Photo_2026-06-08-12-07-58.png", "③ 사진 기록",     "날짜별 진단 누적"),
    (f"{PIC}/KakaoTalk_Photo_2026-06-08-12-09-00.png", "④ 알림",         "위험 시 자동 알림"),
]
ws = Inches(2.95); hs = Inches(5.0)
for i, (path, k, v) in enumerate(shots):
    x = Inches(0.45) + i * Inches(3.15)
    y = Inches(1.05)
    # phone frame
    rounded(s, x - Inches(0.04), y - Inches(0.04), ws + Inches(0.08), hs + Inches(0.08),
            fill=NAVY, line=None)
    add_image_fit(s, path, x, y, ws, hs)
    textbox(s, x, y + hs + Inches(0.12), ws, Inches(0.35),
            k, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    textbox(s, x, y + hs + Inches(0.45), ws, Inches(0.35),
            v, size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ============== 슬라이드 11.5: 앱 사용 시나리오 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "08. 앱 사용 시나리오 — \"위험 감지에서 확인까지\"",
            "농가가 RUST 병해 발생을 발견하는 실제 흐름")

# 왼쪽: 큰 캡처 1장 (구역 현황 위험)
add_image_fit(s, f"{PIC}/KakaoTalk_Photo_2026-06-08-12-08-33.png",
              Inches(0.5), Inches(1.05), Inches(3.0), Inches(5.4))
textbox(s, Inches(0.5), Inches(6.50), Inches(3.0), Inches(0.35),
        "Step 1. 위험 구역 확인", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 중앙: 구역 상세
add_image_fit(s, f"{PIC}/KakaoTalk_Photo_2026-06-08-12-08-42.png",
              Inches(3.85), Inches(1.05), Inches(3.0), Inches(5.4))
textbox(s, Inches(3.85), Inches(6.50), Inches(3.0), Inches(0.35),
        "Step 2. AI 진단 근거 조회", size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# 오른쪽: 알림
add_image_fit(s, f"{PIC}/KakaoTalk_Photo_2026-06-08-12-09-00.png",
              Inches(7.2), Inches(1.05), Inches(3.0), Inches(5.4))
textbox(s, Inches(7.2), Inches(6.50), Inches(3.0), Inches(0.35),
        "Step 3. 알림 수신 및 처리", size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 우측 설명 박스
rounded(s, Inches(10.5), Inches(1.05), Inches(2.6), Inches(5.4), fill=LIGHT, line=BORDER)
textbox(s, Inches(10.65), Inches(1.20), Inches(2.4), Inches(0.4),
        "Storytelling", size=12, bold=True, color=GREEN)
points = [
    ("①", "농가가 앱을 켜면 첫 화면에 정상/위험/주의 카운트가 한 줄 요약으로 표시"),
    ("②", "위험 카드(빨간 테두리)를 탭하면 진단된 잎 사진과 질병명·신뢰도가 모달로 노출"),
    ("③", "동시에 알림 탭에 위험 신호가 누적되어 다른 위치에서도 확인 가능"),
    ("✓", "농가는 별도 도구나 검색 없이 앱 한 번 켜는 것만으로 전 구역 상태 파악"),
]
yy = Inches(1.65)
for ico, txt in points:
    textbox(s, Inches(10.65), yy, Inches(0.25), Inches(0.35),
            ico, size=14, bold=True, color=GREEN)
    textbox(s, Inches(10.95), yy, Inches(2.05), Inches(0.95),
            txt, size=9, color=NAVY)
    yy += Inches(1.05)

# ============== 슬라이드 12: 역할 분담 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "09. 역할 분담", "Team Roles & Responsibilities")
members = [
    ("봉  준  표", "팀장", DGREEN,
     [
        "AI 모델 학습 (EfficientNet-B3)",
        "Android ↔ AI 연동 (PlantDiseaseClassifier.java)",
        "프로젝트 설계 / 일정 총괄",
     ],
     [
        "전처리 (HSV Anomaly Score) · ONNX 변환",
        "그리드 추론 · 오버레이 시각화 로직",
        "ai/ 디렉토리 코드 + README 문서화",
     ]),
    ("이  동  제", "팀원", BLUE,
     [
        "Android 앱 UI / 화면 구성",
        "카메라 촬영 · 이미지 로딩",
     ],
     [
        "메인/촬영/분석/결과 Activity 레이아웃",
        "CameraX 통합 + 갤러리 접근 권한",
        "결과 카드 / 그리드 오버레이 출력",
     ]),
    ("윤  준  영", "팀원", ORANGE,
     [
        "테스트",
        "문서화",
     ],
     [
        "test_images 33장 수집 · 분류",
        "모델 추론 정확도 검증",
        "발표 자료 / README 작성 지원",
     ]),
]
my = Inches(1.15); mh = Inches(5.65); mw = Inches(4.05)
for i, (name, role, color, roles, tasks) in enumerate(members):
    x = Inches(0.7) + i * Inches(4.15)
    rounded(s, x, my, mw, mh, fill=WHITE, line=color, line_w=2)
    rect(s, x, my, mw, Inches(1.05), color)
    textbox(s, x + Inches(0.2), my + Inches(0.15), mw - Inches(0.3), Inches(0.4),
            role, size=12, color=WHITE, bold=True)
    textbox(s, x + Inches(0.2), my + Inches(0.45), mw - Inches(0.3), Inches(0.6),
            name, size=22, bold=True, color=WHITE)
    # 맡은 역할
    textbox(s, x + Inches(0.25), my + Inches(1.18), mw - Inches(0.4), Inches(0.4),
            "맡은 역할", size=12, bold=True, color=color)
    yy = my + Inches(1.50)
    for r in roles:
        textbox(s, x + Inches(0.25), yy, mw - Inches(0.4), Inches(0.4),
                f"•  {r}", size=12, color=NAVY)
        yy += Inches(0.40)
    # 구현 부분
    textbox(s, x + Inches(0.25), yy + Inches(0.15), mw - Inches(0.4), Inches(0.4),
            "구현한 부분", size=12, bold=True, color=color)
    yy += Inches(0.50)
    for t in tasks:
        textbox(s, x + Inches(0.25), yy, mw - Inches(0.4), Inches(0.5),
                f"·  {t}", size=11, color=SLATE)
        yy += Inches(0.42)

# ============== 슬라이드 13: 개발 과정에서 느낀 점 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "10. 개발 과정에서 느낀 점", "Reflections")
reflections = [
    ("📱", "\"앱이 먼저\"라는 관점이 핵심이었다",
     "AI 정확도가 아무리 높아도 사용자가 켜지 않으면 무용지물. 첫 화면에 \"오늘 위험한 구역이 몇 개인가\"가 "
     "한 줄로 보이도록 한 카드뷰 디자인이 앱 전체의 가치를 결정했다.",
     GREEN),
    ("🎨", "Material Design 토큰 시스템의 효용",
     "그린 톤(primary), 빨강(danger), 주황(warning)을 처음에 정한 뒤로는 모든 화면에서 색 고민이 사라졌다. "
     "디자인 시스템을 일찍 정의한 것이 빠른 화면 추가의 비결이었다.",
     BLUE),
    ("🔌", "백엔드 API 응답 스키마가 곧 화면 구조",
     "/analyze가 반환하는 disease_type, confidence가 그대로 카드뷰의 라벨·색상으로 매핑됐다. "
     "백엔드 응답을 먼저 합의한 덕에 앱 개발자는 화면 마크업에만 집중할 수 있었다.",
     ORANGE),
    ("🔁", "실제 캡처 기반 QA의 위력",
     "시뮬레이터가 아닌 실제 안드로이드 단말에서 캡처를 떠보면 토스트 위치, 카드 잘림, 상태바 높이 등 "
     "에뮬레이터에서 못 본 문제가 모두 드러났다. 정기적 실기기 검증이 \"안 깨지는 앱\"을 만든다.",
     DGREEN),
    ("👥", "역할 분담은 인터페이스부터",
     "AI(봉준표), 앱 UI(이동제), QA/문서(윤준영) 셋이 1주차에 PlantDiseaseClassifier 메서드 시그니처와 "
     "응답 JSON 스키마를 먼저 합의한 게 신의 한 수. 이후 각자 병렬로 작업할 수 있었다.",
     RED),
]
yy = Inches(1.15)
for icon, title, body, color in reflections:
    rounded(s, Inches(0.7), yy, Inches(12.0), Inches(1.08), fill=WHITE, line=BORDER)
    rect(s, Inches(0.7), yy, Inches(0.10), Inches(1.08), color)
    textbox(s, Inches(0.9), yy + Inches(0.10), Inches(0.7), Inches(0.85),
            icon, size=28, color=color, align=PP_ALIGN.CENTER)
    textbox(s, Inches(1.7), yy + Inches(0.1), Inches(10.8), Inches(0.4),
            title, size=14, bold=True, color=color)
    textbox(s, Inches(1.7), yy + Inches(0.45), Inches(10.8), Inches(0.6),
            body, size=11, color=NAVY)
    yy += Inches(1.15)

# ============== 슬라이드 14: 향후 개선 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "11. 향후 개선 방향", "What's Next")
future = [
    ("FastAPI 백엔드 완성",
     "현재 README 단계 — 분석 결과 영구 저장 + 구역별 시계열 차트 + 농가 간 비교 통계",
     "Backend", BLUE),
    ("Bluetooth 외부 카메라 연동",
     "비닐하우스 천장 고정 카메라에서 자동 촬영 → 무인 모니터링",
     "Hardware", ORANGE),
    ("실시간 영상 추론",
     "정지 사진이 아닌 CameraX Preview Frame을 30fps로 분석 — 잎을 비추기만 하면 즉시 진단",
     "Realtime", RED),
    ("질병별 방제 DB 구축",
     "23개 질병에 대한 농약/조치/주의사항을 농촌진흥청 데이터와 연계해 자동 추천",
     "Domain", DGREEN),
    ("작물 확장 (현재 5종 → 15종)",
     "딸기, 상추, 오이 등 한국 비닐하우스 주요 작물로 데이터셋 확장 학습",
     "ML", GREEN),
    ("Push 알림 임계치 학습",
     "농가별 사용 패턴을 보고 알림이 너무 잦지 않도록 임계치를 자동 조정",
     "UX", BLUE),
]
yy = Inches(1.15)
for i, (k, v, tag, color) in enumerate(future):
    col = i % 2; row = i // 2
    x = Inches(0.7) + col * Inches(6.15)
    y = Inches(1.15) + row * Inches(1.75)
    rounded(s, x, y, Inches(6.0), Inches(1.65), fill=WHITE, line=BORDER)
    # 태그
    rounded(s, x + Inches(0.25), y + Inches(0.15), Inches(1.2), Inches(0.35),
            fill=color, line=None)
    textbox(s, x + Inches(0.25), y + Inches(0.17), Inches(1.2), Inches(0.3),
            tag, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(1.55), y + Inches(0.16), Inches(4.3), Inches(0.4),
            k, size=14, bold=True, color=NAVY)
    textbox(s, x + Inches(0.25), y + Inches(0.65), Inches(5.6), Inches(1.0),
            v, size=11, color=SLATE)

# ============== 슬라이드 15: 소스코드 위치 ==============
s = add_blank_slide(); set_bg(s, WHITE)
page_header(s, "12. 소스 코드 / 산출물", "Deliverables")
rounded(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.0), fill=NAVY, line=None)
textbox(s, Inches(1.0), Inches(1.35), Inches(11), Inches(0.4),
        "GitHub Repository", size=14, color=GREEN, bold=True)
textbox(s, Inches(1.0), Inches(1.65), Inches(11), Inches(0.5),
        "github.com/bongjunpyo/smartfarm-upgrade", size=20, bold=True, color=WHITE)

tree = [
    ("smartfarm-upgrade/", "프로젝트 루트", DGREEN, True),
    ("├── ai/",                  "AI 모듈 (Python)", BLUE, False),
    ("│   ├── README.md",        "전체 파이프라인 문서", MUTED, False),
    ("│   ├── model/best_crop_model.onnx",  "학습된 모델 (~41 MB)", GREEN, False),
    ("│   ├── preprocessing/",   "HSV 전처리, 학습, GradCAM 5종 스크립트", MUTED, False),
    ("│   └── test_images/",     "테스트 이미지 33장", MUTED, False),
    ("├── android/",             "Android 통합 코드", ORANGE, False),
    ("│   ├── README.md",        "통합 방법 가이드", MUTED, False),
    ("│   └── PlantDiseaseClassifier.java", "추론 클래스 (단일/그리드)", GREEN, False),
    ("├── backend/",             "FastAPI 백엔드 (설계만)", RED, False),
    ("│   └── README.md",        "API 엔드포인트 + DB 스키마", MUTED, False),
    ("├── report_assets/",       "발표 자료 PNG (8장)", BLUE, False),
    ("└── 모바일프로그래밍_최종발표.pptx", "본 발표 자료", DGREEN, False),
]
yy = Inches(2.40)
for path, desc, color, bold in tree:
    textbox(s, Inches(0.9), yy, Inches(6.0), Inches(0.32),
            path, size=11, bold=bold, color=color, font="Menlo")
    textbox(s, Inches(7.2), yy, Inches(5.7), Inches(0.32),
            desc, size=11, color=SLATE)
    yy += Inches(0.30)

# 주요 커밋
rounded(s, Inches(0.7), Inches(6.40), Inches(12), Inches(0.45), fill=LIGHT, line=BORDER)
textbox(s, Inches(0.95), Inches(6.45), Inches(11), Inches(0.4),
        "최근 커밋: 9f462aa  docs: AI 모듈 README 전면 재작성 · 5a7edbb  android+ai: grid inference (10×10)",
        size=10, color=NAVY, font="Menlo")

# ============== 슬라이드 16: Thank You ==============
s = add_blank_slide()
set_bg(s, NAVY)
rect(s, 0, Inches(6.2), SW, Inches(1.3), DGREEN)
rect(s, 0, Inches(6.0), SW, Inches(0.2), GREEN)
textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.5),
        "감 사 합 니 다", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.6),
        "Thank you for your attention.", size=22,
        color=RGBColor(0xa7, 0xf3, 0xd0), align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4),
        "Q & A", size=20, color=WHITE, align=PP_ALIGN.CENTER)
# 팀 이름
textbox(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
        "팀장 봉준표  ·  팀원 이동제  ·  팀원 윤준영",
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
        "2026.06.07  ·  모바일 프로그래밍",
        size=12, color=RGBColor(0xa7, 0xf3, 0xd0), align=PP_ALIGN.CENTER)

# 저장
prs.save(FILE)
print(f"\nSaved: {FILE}")
print(f"Slides: {len(prs.slides)}")
